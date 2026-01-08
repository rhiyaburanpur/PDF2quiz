import re
import os
import json
import time 
import random 
from PyPDF2 import PdfReader
from typing import List, Dict, Any
from io import BytesIO

try:
    from transformers import pipeline
except ImportError:
    def pipeline(task, model):
        print(f"WARNING: Hugging Face 'transformers' library not found. Using stub for {model}.")
        return lambda x: [{"generated_text": "Could not generate content. Install 'transformers' and 'torch'."}]

def get_qa_pipeline():
    """Loads and caches the Hugging Face QA pipeline."""
    # Using a model that performs well for question generation
    # Other options: 't5-small' or 'valhalla/t5-small-qg' (if available)
    try:
        return pipeline(
            "text2text-generation", 
            model="facebook/bart-large-cnn", 
            # Load the tokenizer as well
        )
    except Exception as e:
        print(f"Error loading Hugging Face model: {e}")
        # Return a stub function if model loading fails
        return None 

#multifile extraction utility functions

def extract_text_from_pdf(uploaded_file, selected_pages):
    """Extracts text from selected PDF pages with generic cleaning."""
    if not uploaded_file: return ""
    try:
        uploaded_file.seek(0)
        reader = PdfReader(uploaded_file)
        full_text = []
        page_indices = [p - 1 for p in selected_pages if 0 <= p - 1 < len(reader.pages)]
        
        for page_index in page_indices:
            page = reader.pages[page_index]
            text = page.extract_text()
            if text: full_text.append(text)
        
        raw_text = "\n".join(full_text)
        cleaned_text = re.sub(r'\s+', ' ', raw_text).strip()
        cleaned_text = re.sub(r'\b[A-Za-z0-9]\b', '', cleaned_text).strip()
        cleaned_text = re.sub(r'\.{2,}', '.', cleaned_text)
        cleaned_text = re.sub(r',{2,}', ',', cleaned_text)
        
        return cleaned_text.strip()
        
    except Exception as e:
        print(f"Error during PDF text extraction: {e}")
        return ""

def extract_text_from_docx(uploaded_file):
    """Extracts text from a DOCX file."""
    try:
        import docx as docx_module
    except ImportError:
        print("ERROR: python-docx library not found.")
        return ""
        
    if not uploaded_file: return ""
    try:
        uploaded_file.seek(0)
        document = docx_module.Document(BytesIO(uploaded_file.read()))
        full_text = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error during DOCX text extraction: {e}")
        return ""

def extract_text_from_pptx(uploaded_file):
    """Extracts text from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: python-pptx library not found.")
        return ""
        
    if not uploaded_file: return ""
    try:
        uploaded_file.seek(0)
        prs = Presentation(BytesIO(uploaded_file.read()))
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Error during PPTX text extraction: {e}")
        return ""


def get_text_content(uploaded_file, selected_pages, file_type):
    """Selects the correct extraction method based on file_type."""
    if file_type == 'pdf':
        return extract_text_from_pdf(uploaded_file, selected_pages)
    elif file_type == 'docx':
        return extract_text_from_docx(uploaded_file)
    elif file_type == 'pptx':
        return extract_text_from_pptx(uploaded_file)
    else:
        return ""

def chunk_text(text, chunk_size=3000):
    """Splits text into chunks of roughly chunk_size based on sentence boundaries."""
    if not text: return []
    
    # Split by major sentence delimiters
    sentences = re.split(r'(?<=[.?!])\s+', text)
    
    current_chunk = ""
    chunks = []
    
    for sentence in sentences:
        if len(current_chunk) + len(sentence) + 1 > chunk_size and current_chunk:
            chunks.append(current_chunk.strip())
            current_chunk = sentence + "." 
        else:
            current_chunk += (" " if current_chunk else "") + sentence
            
    if current_chunk:
        chunks.append(current_chunk.strip())
        
    return chunks

def generate_questions_with_hf(text_chunk: str, num_questions: int) -> List[Dict[str, Any]]:
    """
    Uses the Hugging Face model to generate structured MCQ questions from a text chunk.
    This function is designed to run multiple times across different chunks.
    """
    qa_pipeline = get_qa_pipeline()
    if not qa_pipeline:
        return []

    quiz_data = []
    
    #prompt for the model
    system_prompt = (
        "You are an expert quiz generator. Analyze the text provided below and generate a list of "
        f"{num_questions} unique multiple-choice questions (MCQs). "
        "Each question MUST have one correct answer and exactly three plausible distractors. "
        "Output the result as a single JSON array object, where each element is an object with "
        "keys: 'question', 'correct_answer', and 'distractors' (an array of 3 strings). "
        "Ensure the 'question' includes context from the source text."
    )
    
    user_prompt = f"Source Text: \"{text_chunk[:2500]}...\""
    
    full_prompt = system_prompt + "\n\n" + user_prompt
    
    try:
        generated_response = qa_pipeline(
            full_prompt, 
            max_length=1500,
            do_sample=True,
            top_k=50,
            temperature=0.7,
            num_return_sequences=1
        )
        
        raw_text = generated_response[0]['generated_text']
        
        # Try to clean the output to isolate the pure JSON array
        # This handles cases where the model wraps the JSON in markdown or plain text
        json_match = re.search(r'\[.*\]', raw_text, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            
            #basic JSON cleanup 
            json_str = json_str.replace('\n', ' ').replace('\\n', ' ').strip()
            
            quiz_list = json.loads(json_str)
            
            for item in quiz_list:
                options = [item['correct_answer']] + item['distractors']
                
                quiz_data.append({
                    "question": item['question'],
                    "options": options,
                    "correct_answer": item['correct_answer'],
                    "source_snippet": text_chunk[:500] + "..." 
                })
        
    except json.JSONDecodeError as e:
        print(f"JSON Decoding Error (Model did not output valid JSON): {e}")
    except Exception as e:
        print(f"An unexpected error occurred during generation: {e}")
        
    return quiz_data


#post-processing and deduplication

def post_process_quiz_data(quiz_data: List[Dict[str, Any]], target_count: int) -> List[Dict[str, Any]]:
    """
    Performs deduplication, cleaning, and shuffling of options.
    """
    
    seen_questions = set()
    final_questions = []

    for item in quiz_data:
        question_text = item.get('question', '').strip()
        
        normalized_q = re.sub(r'[^\w\s]', '', question_text).lower()
        
        if normalized_q in seen_questions:
            continue
        
        seen_questions.add(normalized_q)
        
        options = item.get('options', [])
        
        if options and isinstance(item.get('correct_answer'), str):
            random.shuffle(options)
            item['options'] = options
                
        # final collection
        final_questions.append(item)
        
        if len(final_questions) >= target_count:
            break
            
    return final_questions

#main orchestration function

def run_question_generation(uploaded_file, selected_pages, difficulty, num_questions, q_type, mcq_type, file_type):
    """
    Main orchestration function: uses Hugging Face model for generation.
    """
    #extraction and generic cleaning
    text_content = get_text_content(uploaded_file, selected_pages, file_type)
    
    if not text_content:
        return []

    #chunking
    chunks = chunk_text(text_content, chunk_size=3000)
    
    all_quiz_data = []
    
    # we generate a total pool of questions slightly larger than requested (e.g., 1.5x) to account for potential duplicates or invalid JSON outputs.
    
    # distributing the generation request across all chunks
    target_generation_per_chunk = (num_questions * 2 + len(chunks) - 1) // len(chunks) if chunks else num_questions
    
    for chunk in chunks:
        # calling the HF model generation function
        quiz_data_chunk = generate_questions_with_hf(chunk, target_generation_per_chunk)
        all_quiz_data.extend(quiz_data_chunk)
        
    #post-processing/deduplication
    final_data = post_process_quiz_data(all_quiz_data, num_questions)
    
    if not final_data:
        # Fallback if LLM failed to generate or parse
        return [{"question": "AI generation failed. Please check console for errors or install required packages.", "options": ["A", "B", "C", "D"], "correct_answer": "A", "source_snippet": "AI Model Error"}]
        
    return final_data
